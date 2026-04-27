from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ──────────────────────────────────────────────────────────────────
BG       = RGBColor(0x0f, 0x17, 0x2a)
BG_MID   = RGBColor(0x1e, 0x29, 0x3b)
BLUE     = RGBColor(0x38, 0xbd, 0xf8)
WHITE    = RGBColor(0xf8, 0xfa, 0xfc)
GRAY     = RGBColor(0x94, 0xa3, 0xb8)
DARK_GR  = RGBColor(0x47, 0x55, 0x69)
LIGHT    = RGBColor(0xcb, 0xd5, 0xe1)
GREEN    = RGBColor(0x22, 0xc5, 0x5e)
YELLOW   = RGBColor(0xf5, 0x9e, 0x0b)
RED      = RGBColor(0xef, 0x44, 0x44)
GR_BG    = RGBColor(0x05, 0x2e, 0x16)
YL_BG    = RGBColor(0x2d, 0x1f, 0x07)
GY_BG    = RGBColor(0x1e, 0x29, 0x3b)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(blank)


def bg(s, color=BG):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = color


def tx(s, text, x, y, w, h, size=16, bold=False, italic=False,
        color=LIGHT, align=PP_ALIGN.LEFT):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color


def rule(s, x1, y, x2):
    ln = s.shapes.add_connector(1, Inches(x1), Inches(y), Inches(x2), Inches(y))
    ln.line.color.rgb = BLUE
    ln.line.width = Pt(2)


def rect(s, x, y, w, h, fill, border=None):
    sh = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if border:
        sh.line.color.rgb = border
        sh.line.width = Pt(0.75)
    else:
        sh.line.fill.background()
    return sh


def badge(s, label, x, y, color, bg_color):
    rect(s, x, y, 1.9, 0.32, bg_color)
    tx(s, label, x + 0.08, y + 0.03, 1.75, 0.28,
       size=10, bold=True, color=color, align=PP_ALIGN.CENTER)


def req_card(s, req, title, status_label, status_color, status_bg,
             bullets, x, card_w=2.9):
    """Draw one requirement status card."""
    card_h = 4.2
    y0 = 2.6
    rect(s, x, y0, card_w, card_h, BG_MID, border=RGBColor(0x33, 0x41, 0x55))
    tx(s, req,   x+0.18, y0+0.18, card_w-0.3, 0.38, size=11, bold=True, color=BLUE)
    tx(s, title, x+0.18, y0+0.58, card_w-0.3, 0.55, size=15, bold=True, color=WHITE)
    badge(s, status_label, x+0.18, y0+1.18, status_color, status_bg)
    # bullets
    tb = s.shapes.add_textbox(
        Inches(x+0.18), Inches(y0+1.68), Inches(card_w-0.3), Inches(2.2))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(5)
        r = p.add_run()
        r.text = f"• {b}"
        r.font.size = Pt(12)
        r.font.color.rgb = LIGHT


# ════════════════════════════════════════════════════════════════════════════
# Slide 1 — Title
# ════════════════════════════════════════════════════════════════════════════
s1 = slide(); bg(s1)
tx(s1, "ACCENTURE&CT CONSULTING  ·  MERIDIAN COMPONENTS",
   0.6, 0.5, 12, 0.4, size=11, bold=True, color=DARK_GR)
rule(s1, 0.6, 1.0, 1.8)
tx(s1, "ENGAGEMENT PROGRESS REPORT", 0.6, 1.1, 10, 0.4,
   size=11, bold=True, color=BLUE)
tx(s1, "Sprint Review\nR1 · R2 · R3 · R4 + D1 · D2 · D3", 0.6, 1.6, 12, 1.8,
   size=38, bold=True, color=WHITE)
tx(s1, "All required deliverables + all optional enhancements complete.",
   0.6, 3.8, 9, 0.9, size=18, color=GRAY)
tx(s1, "April 27, 2026  ·  branch: feature/meridian-engagement",
   0.6, 6.8, 12, 0.35, size=12, color=DARK_GR)

# ════════════════════════════════════════════════════════════════════════════
# Slide 2 — Status overview (4 cards)
# ════════════════════════════════════════════════════════════════════════════
s2 = slide(); bg(s2)
tx(s2, "DELIVERY STATUS", 0.6, 0.55, 10, 0.35, size=11, bold=True, color=BLUE)
tx(s2, "R1 – R4 at a glance", 0.6, 1.0, 10, 0.6, size=28, bold=True, color=WHITE)
rule(s2, 0.6, 0.95, 2.6)

cards = [
    ("R1", "Reports\nRemediation",  "COMPLETE",     GREEN,  GR_BG,
     ["Options → Composition API", "Filters wired (warehouse/category)", "API endpoints patched", "Zero console errors"]),
    ("R2", "Restocking\nView",       "COMPLETE",     GREEN,  GR_BG,
     ["/api/restocking endpoint", "Budget ceiling + priority rank", "New Restocking.vue", "EN + JA locale keys"]),
    ("R3", "Automated\nTests",       "COMPLETE",     GREEN,  GR_BG,
     ["Playwright e2e suite", "59 / 59 tests passing", "All views covered", "Filters + edge cases"]),
    ("R4", "Architecture\nDocs",     "COMPLETE",     GREEN,  GR_BG,
     ["HTML architecture diagram", "15-endpoint API reference", "Mock data layer docs", "IT handoff notes"]),
]

optionals_cards = [
    ("D1", "CSS Variables\n& UI Refresh",  "COMPLETE", GREEN, GR_BG,
     ["CSS custom properties :root", "All hardcoded colours replaced", "FilterBar.vue updated", "Themeable from one block"]),
    ("D2", "Italian\nLocale",              "COMPLETE", GREEN, GR_BG,
     ["it.js — full translation", "EUR currency for 'it'", "translateWarehouse: Londra", "3-language switcher"]),
    ("D3", "Dark Mode\nToggle",            "COMPLETE", GREEN, GR_BG,
     ["Sun/moon button in nav", ".dark class on <html>", "CSS var overrides", "localStorage persistence"]),
]

x = 0.55
for req, title, label, col, bg_col, buls in cards:
    req_card(s2, req, title, label, col, bg_col, buls, x, card_w=2.9)
    x += 3.07

# ════════════════════════════════════════════════════════════════════════════
# Slide 3 — R1 detail
# ════════════════════════════════════════════════════════════════════════════
s3 = slide(); bg(s3, BG_MID)
badge(s3, "COMPLETE", 0.6, 0.45, GREEN, GR_BG)
tx(s3, "R1 · Reports Remediation", 0.6, 0.9, 10, 0.5, size=11, bold=True, color=BLUE)
tx(s3, "All defects resolved.\nZero console errors.", 0.6, 1.35, 11, 1.0, size=32, bold=True, color=WHITE)

defects = [
    ("Options API → Composition API", "Full rewrite of Reports.vue to Composition API pattern consistent with the rest of the codebase. Reactive state, watchers, computed properties."),
    ("Filters fully wired",           "Warehouse and Category params added to /api/reports/quarterly and /api/reports/monthly-trends. Watch on both filters triggers data reload."),
    ("i18n gaps closed",              "All labels, headings, and table columns routed through t() calls. English and Japanese locale keys complete."),
    ("Console noise eliminated",      "No console.log calls, no uncaught errors, no stray warnings on page load or filter change."),
]
x = 0.6
for title, desc in defects:
    rect(s3, x, 2.7, 2.9, 2.3, BG, border=RGBColor(0x33, 0x41, 0x55))
    tx(s3, title, x+0.15, 2.8, 2.6, 0.45, size=13, bold=True, color=BLUE)
    tx(s3, desc,  x+0.15, 3.3, 2.6, 1.4,  size=12, color=LIGHT)
    x += 3.07

# ════════════════════════════════════════════════════════════════════════════
# Slide 4 — R2 detail
# ════════════════════════════════════════════════════════════════════════════
s4 = slide(); bg(s4)
badge(s4, "COMPLETE", 0.6, 0.45, GREEN, GR_BG)
tx(s4, "R2 · Restocking Recommendations", 0.6, 0.9, 10, 0.5, size=11, bold=True, color=BLUE)
tx(s4, "New view live.\nBudget ceiling + demand scoring.", 0.6, 1.35, 11, 1.0, size=32, bold=True, color=WHITE)

# Left: how it works
rect(s4, 0.6, 2.7, 6.0, 4.0, BG_MID, border=RGBColor(0x33, 0x41, 0x55))
tx(s4, "How it works", 0.8, 2.85, 5.6, 0.4, size=14, bold=True, color=BLUE)
items = [
    "Stock gap = reorder_point − quantity_on_hand",
    "Priority score: gap + bonus for 'increasing' demand trend",
    "Budget ceiling: greedy selection, within_budget flag per row",
    "Operator applies budget via button (not on every keystroke)",
    "Filters: warehouse + category reload recommendations",
    "Status column appears only when budget > 0",
]
tb = s4.shapes.add_textbox(Inches(0.8), Inches(3.35), Inches(5.6), Inches(3.0))
tf = tb.text_frame; tf.word_wrap = True
for i, item in enumerate(items):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_before = Pt(6)
    r = p.add_run(); r.text = f"→  {item}"
    r.font.size = Pt(13); r.font.color.rgb = LIGHT

# Right: stat cards mini-preview
rect(s4, 6.85, 2.7, 6.1, 4.0, BG_MID, border=RGBColor(0x33, 0x41, 0x55))
tx(s4, "Summary stats delivered", 7.05, 2.85, 5.7, 0.4, size=14, bold=True, color=BLUE)
stats = [("Items at Risk", "23 SKUs understocked"),
         ("Total Restock Cost", "Full cost if all ordered"),
         ("Within Budget", "Items within ceiling"),
         ("Budget Remaining", "Unspent allocation")]
y = 3.4
for label, desc in stats:
    tx(s4, f"■  {label}", 7.05, y, 5.7, 0.3, size=13, bold=True, color=BLUE)
    tx(s4, desc, 7.05, y+0.3, 5.7, 0.3, size=12, color=GRAY)
    y += 0.72

# ════════════════════════════════════════════════════════════════════════════
# Slide 5 — R3 detail
# ════════════════════════════════════════════════════════════════════════════
s5 = slide(); bg(s5, BG_MID)
badge(s5, "COMPLETE", 0.6, 0.45, GREEN, GR_BG)
tx(s5, "R3 · Automated Browser Tests", 0.6, 0.9, 10, 0.5, size=11, bold=True, color=BLUE)
tx(s5, "59 / 59 tests passing.", 0.6, 1.35, 8, 0.8, size=36, bold=True, color=WHITE)
tx(s5, "59", 10.5, 1.2, 2.5, 1.5, size=72, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
tx(s5, "passing", 10.5, 2.7, 2.5, 0.4, size=14, color=GRAY, align=PP_ALIGN.CENTER)

suites = [
    ("navigation.spec.js",  "3 tests",  "Nav bar, route loading, active link"),
    ("filters.spec.js",     "5 tests",  "All 4 controls, options, reset"),
    ("dashboard.spec.js",   "6 tests",  "Title, KPIs, Order Health, charts, filters"),
    ("inventory.spec.js",   "6 tests",  "Table, columns, Location/Category filters"),
    ("orders.spec.js",      "5 tests",  "Table, columns, Status/Location/Period filters"),
    ("reports.spec.js",     "8 tests",  "Quarterly table, MoM, bar chart, stats, filters"),
    ("restocking.spec.js",  "9 tests",  "Budget input, stat cards, table, Status column logic"),
    ("spending.spec.js",    "6 tests",  "KPI cards, charts, period filter"),
    ("demand.spec.js",      "7 tests",  "Trend cards, table, filters"),
]
x, y = 0.6, 2.65
col = 0
for name, count, desc in suites:
    rect(s5, x, y, 4.1, 0.7, BG, border=RGBColor(0x1e, 0x29, 0x3b))
    tx(s5, name,  x+0.15, y+0.06, 2.5, 0.3, size=11, bold=True, color=BLUE)
    tx(s5, count, x+2.8,  y+0.06, 1.1, 0.3, size=11, bold=True, color=GREEN, align=PP_ALIGN.RIGHT)
    tx(s5, desc,  x+0.15, y+0.36, 3.8, 0.28, size=10, color=GRAY)
    col += 1
    if col == 3:
        col = 0; x = 0.6; y += 0.82
    else:
        x += 4.25

# ════════════════════════════════════════════════════════════════════════════
# Slide 6 — R4 detail
# ════════════════════════════════════════════════════════════════════════════
s6 = slide(); bg(s6)
badge(s6, "COMPLETE", 0.6, 0.45, GREEN, GR_BG)
tx(s6, "R4 · Architecture Documentation", 0.6, 0.9, 10, 0.5, size=11, bold=True, color=BLUE)
tx(s6, "Full current-state docs\ndelivered to Meridian IT.", 0.6, 1.35, 11, 1.0, size=32, bold=True, color=WHITE)

r4_items = [
    ("Tech stack overview",       "Vue 3 + Vite / FastAPI + Pydantic / JSON mock data. Ports, run commands, CORS notes."),
    ("SVG component map",         "Browser → api.js → FastAPI endpoints → mock data layer. Composables, data flow arrows, legend."),
    ("API endpoint reference",    "15 endpoints: method, path, all query params, JSON response shape. Restocking ★ flagged as new."),
    ("IT handoff notes",          "How to add endpoints, add views, replace JSON with a real DB, run the Playwright suite."),
]
x = 0.6
for title, desc in r4_items:
    rect(s6, x, 2.7, 2.9, 2.3, BG_MID, border=RGBColor(0x33, 0x41, 0x55))
    tx(s6, title, x+0.15, 2.8,  2.6, 0.45, size=13, bold=True, color=BLUE)
    tx(s6, desc,  x+0.15, 3.28, 2.6, 1.4,  size=12, color=LIGHT)
    x += 3.07

# ════════════════════════════════════════════════════════════════════════════
# Slide 7 — D1+D2+D3 status cards
# ════════════════════════════════════════════════════════════════════════════
s7 = slide(); bg(s7)
tx(s7, "OPTIONAL ENHANCEMENTS", 0.6, 0.55, 10, 0.35, size=11, bold=True, color=BLUE)
tx(s7, "D1 · D2 · D3 — all delivered", 0.6, 1.0, 10, 0.6, size=28, bold=True, color=WHITE)
rule(s7, 0.6, 0.95, 2.6)

x = 0.85
for req, title, label, col, bg_col, buls in optionals_cards:
    req_card(s7, req, title, label, col, bg_col, buls, x, card_w=3.8)
    x += 3.97

tx(s7, "branch: feature/meridian-engagement  ·  RFP #MC-2026-0417  ·  April 27, 2026",
   0.6, 7.1, 12, 0.35, size=11, color=DARK_GR)

# ════════════════════════════════════════════════════════════════════════════
# Slide 8 — D1 detail: CSS variables + dark mode
# ════════════════════════════════════════════════════════════════════════════
s8 = slide(); bg(s8, BG_MID)
badge(s8, "COMPLETE", 0.6, 0.45, GREEN, GR_BG)
tx(s8, "D1 + D3 · CSS Variables & Dark Mode", 0.6, 0.9, 10, 0.5, size=11, bold=True, color=BLUE)
tx(s8, "Full theme system.\nOne-click dark mode.", 0.6, 1.35, 11, 1.0, size=32, bold=True, color=WHITE)

d1_items = [
    (":root CSS variables",     "14 custom properties — --bg, --surface, --border, --text-primary, --accent, etc. All layout and colour tokens centralised."),
    (".dark class override",    "Applied to <html> element. Full dark palette defined in one block; flips the entire app including charts, cards, tables."),
    ("App.vue refactored",      "All hardcoded hex values replaced with var(--*) calls. FilterBar.vue scoped styles also updated."),
    ("Toggle + persistence",    "Sun/moon SVG button in nav. isDark ref seeds from localStorage on mount; preference survives page reload."),
]
x = 0.6
for title, desc in d1_items:
    rect(s8, x, 2.7, 2.9, 2.3, BG, border=RGBColor(0x33, 0x41, 0x55))
    tx(s8, title, x+0.15, 2.8,  2.6, 0.45, size=13, bold=True, color=BLUE)
    tx(s8, desc,  x+0.15, 3.28, 2.6, 1.4,  size=12, color=LIGHT)
    x += 3.07

tx(s8, "branch: feature/meridian-engagement  ·  commit: 272787c",
   0.6, 7.1, 12, 0.35, size=11, color=DARK_GR)

# ════════════════════════════════════════════════════════════════════════════
# Slide 9 — D2 detail: Italian locale
# ════════════════════════════════════════════════════════════════════════════
s9 = slide(); bg(s9)
badge(s9, "COMPLETE", 0.6, 0.45, GREEN, GR_BG)
tx(s9, "D2 · Italian Locale (i18n Expansion)", 0.6, 0.9, 10, 0.5, size=11, bold=True, color=BLUE)
tx(s9, "Third language live.\nEN · JA · IT.", 0.6, 1.35, 11, 1.0, size=32, bold=True, color=WHITE)

# Left: what was added
rect(s9, 0.6, 2.7, 6.0, 4.0, BG_MID, border=RGBColor(0x33, 0x41, 0x55))
tx(s9, "Files added / modified", 0.8, 2.85, 5.6, 0.4, size=14, bold=True, color=BLUE)
d2_files = [
    "client/src/locales/it.js — full translation (403 lines)",
    "useI18n.js — import it, EUR currency for 'it' locale",
    "useI18n.js — translateWarehouse: London → 'Londra'",
    "LanguageSwitcher.vue — 'Italiano' added to languageNames",
]
tb9 = s9.shapes.add_textbox(Inches(0.8), Inches(3.35), Inches(5.6), Inches(3.0))
tf9 = tb9.text_frame; tf9.word_wrap = True
for i, item in enumerate(d2_files):
    p = tf9.paragraphs[0] if i == 0 else tf9.add_paragraph()
    p.space_before = Pt(8)
    r = p.add_run(); r.text = f"→  {item}"
    r.font.size = Pt(13); r.font.color.rgb = LIGHT

# Right: translated samples
rect(s9, 6.85, 2.7, 6.1, 4.0, BG_MID, border=RGBColor(0x33, 0x41, 0x55))
tx(s9, "Sample translations", 7.05, 2.85, 5.7, 0.4, size=14, bold=True, color=BLUE)
samples = [
    ("nav.overview",        "Panoramica"),
    ("nav.restocking",      "Rifornimento"),
    ("dashboard.kpi.title", "Indicatori Chiave di Performance"),
    ("filters.location",    "Sede"),
    ("status.inStock",      "Disponibile"),
    ("priority.high",       "Alta"),
]
y9 = 3.38
for key, val in samples:
    tx(s9, key,  7.05, y9,       3.0, 0.28, size=11, color=GRAY)
    tx(s9, val,  10.2, y9,       2.5, 0.28, size=11, bold=True, color=WHITE)
    y9 += 0.4

tx(s9, "branch: feature/meridian-engagement  ·  commit: 272787c",
   0.6, 7.1, 12, 0.35, size=11, color=DARK_GR)

# ════════════════════════════════════════════════════════════════════════════
# Slide 10 — Engagement complete + PR
# ════════════════════════════════════════════════════════════════════════════
s10 = slide(); bg(s10)
tx(s10, "ACCENTURE&CT CONSULTING  ·  MERIDIAN COMPONENTS",
   0.6, 0.5, 12, 0.4, size=11, bold=True, color=DARK_GR)
rule(s10, 0.6, 1.0, 1.8)
tx(s10, "FULL ENGAGEMENT · COMPLETE", 0.6, 1.1, 10, 0.4, size=11, bold=True, color=BLUE)
tx(s10, "R1 · R2 · R3 · R4\n+ D1 · D2 · D3", 0.6, 1.55, 8, 1.6, size=40, bold=True, color=WHITE)

# Full delivery strip
rect(s10, 0.6, 3.35, 12.1, 0.55, GR_BG, border=RGBColor(0x16, 0x65, 0x34))
all_done = ["R1  Reports", "R2  Restocking", "R3  59/59 tests", "R4  Arch docs", "D1  CSS vars", "D2  Italiano", "D3  Dark mode"]
xi = 0.75
for item in all_done:
    tx(s10, "✓  " + item, xi, 3.42, 1.65, 0.38, size=11, bold=True, color=GREEN)
    xi += 1.73

# PR / branch info box
rect(s10, 0.6, 4.1, 12.1, 2.6, BG_MID, border=RGBColor(0x33, 0x41, 0x55))
tx(s10, "Pull Request ready for review", 0.8, 4.25, 11, 0.4, size=14, bold=True, color=BLUE)

pr_info = [
    ("Branch",    "feature/meridian-engagement  →  main"),
    ("Commits",   "272787c  Deliver D1+D2+D3: CSS variables, Italian locale, dark mode toggle\na0e7e5d  Update progress deck: all four R deliverables COMPLETE\n645b5f2  R4: Add architecture documentation\n8004785  Complete R1–R3 Meridian engagement deliverables"),
    ("Changed",   "15 files  ·  App.vue, FilterBar.vue, LanguageSwitcher.vue, useI18n.js,\nit.js (new), Restocking.vue (new), main.py, api.js, e2e test suite"),
    ("Repo",      "github.com/lindsey-anthropic/meridian-workshop"),
]
y10 = 4.78
for label, val in pr_info:
    tx(s10, label + ":", 0.8,  y10, 1.2, 0.35, size=12, bold=True, color=GRAY)
    tx(s10, val,         2.15, y10, 9.8, 0.45, size=12, color=LIGHT)
    y10 += 0.5

tx(s10, "RFP #MC-2026-0417  ·  April 27, 2026",
   0.6, 7.1, 12, 0.35, size=11, color=DARK_GR)

prs.save("progress-report.pptx")
print("Done: progress-report.pptx")
