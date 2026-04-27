from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# Colors
C_BG      = RGBColor(0x0f, 0x17, 0x2a)   # dark background
C_CARD    = RGBColor(0x1e, 0x29, 0x3b)   # card surface
C_TRUST   = RGBColor(0x25, 0x63, 0xeb)   # blue — trust
C_PAIN    = RGBColor(0xf5, 0x9e, 0x0b)   # amber — empathy
C_RESOLVE = RGBColor(0x10, 0xb9, 0x81)   # green — solution
C_WHITE   = RGBColor(0xf8, 0xfa, 0xfc)
C_MUTED   = RGBColor(0x94, 0xa3, 0xb8)
C_DIM     = RGBColor(0x47, 0x55, 0x69)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]  # completely blank

def add_slide():
    return prs.slides.add_slide(blank_layout)

def bg(slide, color=C_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, x, y, w, h, color, radius=False):
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def label(slide, text, x, y, color=C_TRUST, size=10):
    tb = slide.shapes.add_textbox(x, y, Inches(10), Pt(20))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = text.upper()
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = color
    p.font.name = "Segoe UI"

def heading(slide, text, x, y, w=Inches(11), size=36, color=C_WHITE):
    tb = slide.shapes.add_textbox(x, y, w, Inches(2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = color
    p.font.name = "Segoe UI"

def body(slide, text, x, y, w=Inches(5), h=Inches(3), size=14, color=C_MUTED, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.name = "Segoe UI"
    return tb

def accent_bar(slide, color):
    rect(slide, 0, 0, Inches(0.08), H, color)

def card_box(slide, x, y, w, h, title, subtitle, body_text, accent=C_TRUST):
    r = rect(slide, x, y, w, h, C_CARD)
    # title label
    tb = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.15), w - Inches(0.4), Inches(0.3))
    p = tb.text_frame.paragraphs[0]
    p.text = title.upper()
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = accent
    p.font.name = "Segoe UI"
    # subtitle
    tb2 = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.5), w - Inches(0.4), Inches(0.4))
    p2 = tb2.text_frame.paragraphs[0]
    p2.text = subtitle
    p2.font.size = Pt(13)
    p2.font.bold = True
    p2.font.color.rgb = C_WHITE
    p2.font.name = "Segoe UI"
    # body
    tb3 = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.95), w - Inches(0.4), h - Inches(1.1))
    tb3.text_frame.word_wrap = True
    p3 = tb3.text_frame.paragraphs[0]
    p3.text = body_text
    p3.font.size = Pt(11)
    p3.font.color.rgb = C_MUTED
    p3.font.name = "Segoe UI"

def highlight_box(slide, text, x, y, w, h, accent=C_TRUST):
    rect(slide, x, y, w, h, C_CARD)
    rect(slide, x, y, Inches(0.06), h, accent)
    tb = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.15), w - Inches(0.35), h - Inches(0.3))
    tb.text_frame.word_wrap = True
    p = tb.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(13)
    p.font.color.rgb = RGBColor(0xcb, 0xd5, 0xe1)
    p.font.name = "Segoe UI"

pad_x = Inches(0.6)
pad_y = Inches(0.7)

# ─── SLIDE 1 — Title ─────────────────────────────────────────────────────────
s = add_slide(); bg(s); accent_bar(s, C_TRUST)
label(s, "RFP MC-2026-0417 · Risposta Accenture", pad_x, Inches(2.2))
heading(s, "Modernizzazione\nInventory Dashboard", pad_x, Inches(2.6), size=44)
body(s, "Proposta a Meridian Components, Inc.  —  Aprile 2026", pad_x, Inches(5.1), size=16, color=C_MUTED)

# ─── SLIDE 2 — Il problema ────────────────────────────────────────────────────
s = add_slide(); bg(s); accent_bar(s, C_PAIN)
label(s, "Contesto", pad_x, pad_y, color=C_PAIN)
heading(s, "Il sistema funziona. Ma non abbastanza.", pad_x, Inches(1.2), size=30)
cw = Inches(3.8); ch = Inches(3.5); cy = Inches(2.2)
card_box(s, pad_x,              cy, cw, ch, "PROBLEMA 1", "Reports difettoso",
         "Filtri non cablati, i18n incompleta, pattern API inconsistenti. Il vendor precedente ha chiuso senza risolvere.", C_PAIN)
card_box(s, pad_x + cw + Inches(0.2), cy, cw, ch, "PROBLEMA 2", "Zero test automatizzati",
         "IT non approva modifiche. Il sistema è bloccato: nessuna feature nuova può andare in produzione in sicurezza.", C_PAIN)
card_box(s, pad_x + (cw + Inches(0.2))*2, cy, cw, ch, "PROBLEMA 3", "Manca il Restocking",
         "Il team operativo gestisce gli ordini a mano. È la funzionalità più attesa.", C_PAIN)

# ─── SLIDE 3 — Analisi ───────────────────────────────────────────────────────
s = add_slide(); bg(s); accent_bar(s, C_PAIN)
label(s, "Analisi", pad_x, pad_y, color=C_PAIN)
heading(s, "Come leggiamo la situazione", pad_x, Inches(1.2), size=30)
bullets = [
    "IT è il gatekeeper: senza test (R3), niente va avanti.",
    "Tanaka è la champion: il Restocking è la sua priorità operativa da mesi.",
    "Okafor vuole prevedibilità: timeline e not-to-exceed > prezzo basso.",
    "Il codice esistente è il punto di partenza — l'handoff del vendor era minimale.",
]
for i, b in enumerate(bullets):
    body(s, f"—  {b}", pad_x, Inches(2.3) + Inches(0.65)*i, w=Inches(7), size=14)
highlight_box(s, "Sequenza giusta: sbloccare (R3+R1) → estendere (R2) → documentare (R4).",
              Inches(8.5), Inches(2.5), Inches(4.2), Inches(1.2), C_PAIN)

# ─── SLIDE 4 — R1 ────────────────────────────────────────────────────────────
s = add_slide(); bg(s); accent_bar(s, C_RESOLVE)
label(s, "R1 — Obbligatorio", pad_x, pad_y, color=C_RESOLVE)
heading(s, "Remediation Reports", pad_x, Inches(1.2), size=32)
for i, b in enumerate([
    "Discovery autonoma sul codice esistente",
    "Audit filtri, i18n, pattern API, console noise",
    "Fix incrementali con verifica manuale",
    "Completato prima della scrittura dei test (R3)",
]):
    body(s, f"—  {b}", pad_x, Inches(2.4) + Inches(0.6)*i, w=Inches(7), size=14, color=C_MUTED)
card_box(s, Inches(8.5), Inches(2.2), Inches(4.2), Inches(2.2),
         "ASSUNZIONE", "Nessun elenco bug da Meridian",
         "La discovery è interamente a nostro carico. Stimiamo con margine conservativo.", C_RESOLVE)

# ─── SLIDE 5 — R2 ────────────────────────────────────────────────────────────
s = add_slide(); bg(s); accent_bar(s, C_RESOLVE)
label(s, "R2 — Obbligatorio", pad_x, pad_y, color=C_RESOLVE)
heading(s, "Vista Restocking", pad_x, Inches(1.2), size=32)
cw = Inches(3.8); ch = Inches(3.0); cy = Inches(2.5)
card_box(s, pad_x, cy, cw, ch, "INPUT", "Stock + Domanda + Budget",
         "L'operatore inserisce un budget ceiling. Il sistema incrocia stock attuale e forecast di domanda.", C_RESOLVE)
card_box(s, pad_x + cw + Inches(0.2), cy, cw, ch, "OUTPUT", "Ordini raccomandati",
         "Tabella per magazzino e categoria, con quantità e costo stimato. Modificabile prima della conferma.", C_RESOLVE)
card_box(s, pad_x + (cw+Inches(0.2))*2, cy, cw, ch, "STACK", "Vue 3 + FastAPI",
         "Nuovo endpoint /api/restocking. UX concordata con Tanaka prima dell'implementazione.", C_RESOLVE)

# ─── SLIDE 6 — R3 + R4 ───────────────────────────────────────────────────────
s = add_slide(); bg(s); accent_bar(s, C_RESOLVE)
label(s, "R3 + R4 — Obbligatori", pad_x, pad_y, color=C_RESOLVE)
heading(s, "Test automatizzati e documentazione", pad_x, Inches(1.2), size=28)
card_box(s, pad_x, Inches(2.3), Inches(5.8), Inches(2.5),
         "R3 — PLAYWRIGHT", "Suite di test browser",
         "Copertura sui flussi critici: Dashboard, Reports con filtri, Restocking. Eseguibile in CI. Scritto dopo R1 e R2.", C_RESOLVE)
card_box(s, pad_x + Inches(6.0), Inches(2.3), Inches(5.8), Inches(2.5),
         "R4 — ARCHITETTURA", "Documentazione per IT",
         "Generata dal codice attuale. Formato HTML standalone: mappa componenti, schema API, flusso dati.", C_RESOLVE)
highlight_box(s, "R3 sblocca IT per approvare modifiche future. È non-negoziabile anche se formalmente è uno dei quattro required.",
              pad_x, Inches(5.1), Inches(11.8), Inches(1.2), C_RESOLVE)

# ─── SLIDE 7 — Timeline ──────────────────────────────────────────────────────
s = add_slide(); bg(s); accent_bar(s, C_TRUST)
label(s, "Piano di consegna", pad_x, pad_y)
heading(s, "8 settimane, 4 fasi", pad_x, Inches(1.2), size=30)
phases = [
    ("FASE 1", "Onboarding + R4", "Discovery sul codice, documentazione architetturale", "Sett. 1–2", 0.25),
    ("FASE 2", "R1 — Reports",    "Audit e remediation dei difetti",                    "Sett. 2–4", 0.50),
    ("FASE 3", "R2 — Restocking", "Design UX + implementazione frontend/backend",       "Sett. 4–7", 0.75),
    ("FASE 4", "R3 — Test + Go-live", "Suite Playwright, rilascio, approvazione IT",   "Sett. 7–8", 1.00),
]
for i, (lbl, title, desc, weeks, pct) in enumerate(phases):
    y = Inches(2.3) + Inches(1.05)*i
    ph = Inches(0.9)
    rect(s, pad_x, y, Inches(11.8), ph, C_CARD)
    # phase label
    tb = s.shapes.add_textbox(pad_x + Inches(0.15), y + Inches(0.1), Inches(1.2), Inches(0.3))
    p = tb.text_frame.paragraphs[0]; p.text = lbl; p.font.size = Pt(9); p.font.bold = True
    p.font.color.rgb = C_TRUST; p.font.name = "Segoe UI"
    # title
    tb2 = s.shapes.add_textbox(pad_x + Inches(1.5), y + Inches(0.06), Inches(4), Inches(0.35))
    p2 = tb2.text_frame.paragraphs[0]; p2.text = title; p2.font.size = Pt(13); p2.font.bold = True
    p2.font.color.rgb = C_WHITE; p2.font.name = "Segoe UI"
    # desc
    tb3 = s.shapes.add_textbox(pad_x + Inches(1.5), y + Inches(0.45), Inches(4), Inches(0.3))
    p3 = tb3.text_frame.paragraphs[0]; p3.text = desc; p3.font.size = Pt(10)
    p3.font.color.rgb = C_DIM; p3.font.name = "Segoe UI"
    # progress track
    bw = Inches(5.0); bh = Inches(0.08); bx = pad_x + Inches(6.0); by = y + ph/2 - bh/2
    rect(s, bx, by, bw, bh, RGBColor(0x0f, 0x17, 0x2a))
    rect(s, bx, by, int(bw * pct), bh, C_TRUST)
    # weeks
    tb4 = s.shapes.add_textbox(bx + bw + Inches(0.15), y + Inches(0.25), Inches(1.1), Inches(0.35))
    p4 = tb4.text_frame.paragraphs[0]; p4.text = weeks; p4.font.size = Pt(10)
    p4.font.color.rgb = C_DIM; p4.font.name = "Segoe UI"

# ─── SLIDE 8 — Pricing ───────────────────────────────────────────────────────
s = add_slide(); bg(s); accent_bar(s, C_TRUST)
label(s, "Pricing", pad_x, pad_y)
heading(s, "T&M con not-to-exceed", pad_x, Inches(1.2), size=30)
rows = [
    ("Fase 1", "Onboarding + R4",     "10 gg", "€12.000",  C_WHITE,   C_MUTED),
    ("Fase 2", "R1 Reports",           "20 gg", "€21.000",  C_WHITE,   C_MUTED),
    ("Fase 3", "R2 Restocking",        "30 gg", "€31.500",  C_WHITE,   C_MUTED),
    ("Fase 4", "R3 Test + rilascio",   "5 gg",  "€6.000",   C_WHITE,   C_MUTED),
    ("Contingency 15%", "Discovery, iterazioni UX", "—", "€10.500", C_PAIN, C_MUTED),
    ("Not-to-exceed",   "",            "65 gg", "€81.000",  C_WHITE,   C_WHITE),
]
cols = [Inches(1.5), Inches(4.5), Inches(1.5), Inches(1.8)]
xs   = [pad_x, pad_x+Inches(1.6), pad_x+Inches(7.0), pad_x+Inches(9.0)]
for i, (c1, c2, c3, c4, col1, col4) in enumerate(rows):
    y = Inches(2.4) + Inches(0.72)*i
    bg_c = RGBColor(0x1e,0x29,0x3b) if i % 2 == 0 else C_BG
    if i == 5: bg_c = RGBColor(0x0d,0x1f,0x3c)
    rect(s, pad_x, y, Inches(11.8), Inches(0.65), bg_c)
    for txt, x, w, col in [(c1, xs[0], cols[0], col1), (c2, xs[1], cols[1], C_MUTED),
                            (c3, xs[2], cols[2], C_MUTED), (c4, xs[3], cols[3], col4)]:
        tb = s.shapes.add_textbox(x + Inches(0.1), y + Inches(0.15), w, Inches(0.35))
        p = tb.text_frame.paragraphs[0]; p.text = txt
        p.font.size = Pt(14) if i == 5 else Pt(12)
        p.font.bold = (i == 5)
        p.font.color.rgb = col; p.font.name = "Segoe UI"

# ─── SLIDE 9 — Esperienza ────────────────────────────────────────────────────
s = add_slide(); bg(s); accent_bar(s, C_RESOLVE)
label(s, "Esperienza rilevante", pad_x, pad_y, color=C_RESOLVE)
heading(s, "Tre engagement analoghi", pad_x, Inches(1.2), size=30)
cw = Inches(3.8); cy = Inches(2.3); ch = Inches(3.5)
card_box(s, pad_x, cy, cw, ch, "ENGAGEMENT 1 · EMEA",
         "Dashboard operativa manifatturiera",
         "Rilevato e completato sistema incompleto da vendor precedente. Vue 3 + FastAPI.\n→ R1 + R2", C_RESOLVE)
card_box(s, pad_x + cw + Inches(0.2), cy, cw, ch, "ENGAGEMENT 2 · SaaS",
         "Test automation logistics",
         "Suite Playwright su frontend Vue 3 privo di test. IT come gatekeeper.\n→ R3", C_RESOLVE)
card_box(s, pad_x + (cw+Inches(0.2))*2, cy, cw, ch, "ENGAGEMENT 3 · APAC",
         "i18n + UI refresh warehouse",
         "Estensione i18n a team giapponese e refresh design system.\n→ D1 + D2", C_RESOLVE)

# ─── SLIDE 10 — Desiderati ───────────────────────────────────────────────────
s = add_slide(); bg(s); accent_bar(s, C_TRUST)
label(s, "D1 · D2 · D3 — Desiderati", pad_x, pad_y)
heading(s, "Pronti quando il perimetro lo permette", pad_x, Inches(1.2), size=28)
cw = Inches(3.8); cy = Inches(2.3); ch = Inches(2.8)
card_box(s, pad_x, cy, cw, ch, "D1", "UI modernization",
         "Refresh visivo con design tokens. Direzione concordata con il team operativo, non imposta.", C_TRUST)
card_box(s, pad_x + cw + Inches(0.2), cy, cw, ch, "D2", "Internazionalizzazione",
         "Estensione i18n al team di Tokyo (giapponese). Esperienza diretta su contesti simili.", C_TRUST)
card_box(s, pad_x + (cw+Inches(0.2))*2, cy, cw, ch, "D3", "Dark mode",
         "Toggle lato operatore con CSS custom properties. Per le stazioni in ambienti a bassa luminosità.", C_TRUST)
highlight_box(s, "Stima separata disponibile su richiesta. Orientativamente 3–4 settimane aggiuntive.",
              pad_x, Inches(5.5), Inches(11.8), Inches(0.9))

# ─── SLIDE 11 — Perché noi ───────────────────────────────────────────────────
s = add_slide(); bg(s); accent_bar(s, C_RESOLVE)
label(s, "Perché Accenture", pad_x, pad_y, color=C_RESOLVE)
heading(s, "Lo stesso stack.\nLo stesso contesto.", pad_x, Inches(1.1), size=32)
for i, b in enumerate([
    "Già lavorato su sistemi Vue 3 + FastAPI lasciati incompleti da vendor precedenti.",
    "Conosciamo il pattern IT-gatekeeper: R3 non è optional, è il prerequisito di tutto.",
    "Lavoro incrementale: ogni fase ha un deliverable verificabile.",
    "Not-to-exceed: certezza sul budget come impegno contrattuale.",
]):
    body(s, f"—  {b}", pad_x, Inches(2.8) + Inches(0.65)*i, w=Inches(7.5), size=13, color=C_MUTED)

# Big numbers
for val, lbl, x in [("8", "settimane", Inches(9.2)), ("€81K", "not-to-exceed", Inches(11.0))]:
    tb = s.shapes.add_textbox(x, Inches(2.5), Inches(1.8), Inches(1.4))
    p = tb.text_frame.paragraphs[0]; p.text = val; p.font.size = Pt(52); p.font.bold = True
    p.font.color.rgb = C_RESOLVE; p.font.name = "Segoe UI"; p.alignment = PP_ALIGN.CENTER
    tb2 = s.shapes.add_textbox(x, Inches(3.9), Inches(1.8), Inches(0.4))
    p2 = tb2.text_frame.paragraphs[0]; p2.text = lbl; p2.font.size = Pt(11)
    p2.font.color.rgb = C_MUTED; p2.font.name = "Segoe UI"; p2.alignment = PP_ALIGN.CENTER

# ─── SLIDE 12 — Next steps ───────────────────────────────────────────────────
s = add_slide(); bg(s); accent_bar(s, C_RESOLVE)
label(s, "Prossimi passi", pad_x, pad_y, color=C_RESOLVE)
heading(s, "Come andiamo avanti", pad_x, Inches(1.2), size=32)
cw = Inches(3.8); cy = Inches(2.4); ch = Inches(2.8)
card_box(s, pad_x, cy, cw, ch, "STEP 1", "Aggiudicazione",
         "Firma contratto entro la settimana di aggiudicazione. Kick-off entro 3 giorni lavorativi.", C_RESOLVE)
card_box(s, pad_x + cw + Inches(0.2), cy, cw, ch, "STEP 2", "Accesso ambiente",
         "Accesso al repo e all'ambiente di sviluppo entro fine settimana 1. Dipendenza lato Meridian.", C_RESOLVE)
card_box(s, pad_x + (cw+Inches(0.2))*2, cy, cw, ch, "STEP 3", "Sessione con Tanaka",
         "~1 ora prima della Fase 3 per validare la UX del Restocking. Nessuna implementazione prima.", C_RESOLVE)
highlight_box(s, "Domande? procurement@meridiancomponents.example  —  Scadenza chiarimenti: 28 aprile 2026",
              pad_x, Inches(5.5), Inches(11.8), Inches(0.9), C_RESOLVE)

prs.save("proposal/capabilities-deck.pptx")
print("Generato: proposal/capabilities-deck.pptx")
