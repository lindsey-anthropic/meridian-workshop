from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colors ──────────────────────────────────────────────────────────────────
BG       = RGBColor(0x0f, 0x17, 0x2a)   # dark navy
CARD     = RGBColor(0x1e, 0x29, 0x3b)   # card bg
ACCENT   = RGBColor(0x3b, 0x82, 0xf6)   # blue
ACCENT2  = RGBColor(0x63, 0x66, 0xf1)   # indigo
WHITE    = RGBColor(0xf8, 0xfa, 0xfc)
MUTED    = RGBColor(0x94, 0xa3, 0xb8)
DIM      = RGBColor(0x47, 0x55, 0x69)
GREEN    = RGBColor(0x86, 0xef, 0xac)
AMBER    = RGBColor(0xfc, 0xd3, 0x4d)
PURPLE   = RGBColor(0xd8, 0xb4, 0xfe)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]  # completely blank


def add_slide():
    return prs.slides.add_slide(blank_layout)


def bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def accent_bar(slide, color1=ACCENT, color2=ACCENT2):
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.08), H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color1
    bar.line.fill.background()


def box(slide, left, top, width, height, fill_color=CARD, border_color=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    return shape


def txt(slide, text, left, top, width, height,
        size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb


def eyebrow(slide, text, top=Inches(0.6)):
    txt(slide, text.upper(), Inches(0.5), top, Inches(12), Inches(0.4),
        size=10, bold=True, color=ACCENT)


def heading(slide, text, top=Inches(1.0), size=36):
    txt(slide, text, Inches(0.5), top, Inches(12), Inches(1.6),
        size=size, bold=True, color=WHITE)


def body(slide, text, left, top, width, height, size=15, color=MUTED):
    txt(slide, text, left, top, width, height, size=size, color=color)


def tag_pill(slide, text, left, top, fill, text_color):
    pill = slide.shapes.add_shape(1, left, top, Inches(1.4), Inches(0.28))
    pill.fill.solid()
    pill.fill.fore_color.rgb = fill
    pill.line.fill.background()
    tf = pill.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text.upper()
    run.font.size = Pt(9)
    run.font.bold = True
    run.font.color.rgb = text_color


def divider(slide, top):
    line = slide.shapes.add_shape(1, Inches(0.5), top, Inches(12.3), Pt(1))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0x1e, 0x29, 0x3b)
    line.line.fill.background()


def bullet_list(slide, items, left, top, width, size=14, color=MUTED, gap=Inches(0.32)):
    for i, item in enumerate(items):
        dot = slide.shapes.add_shape(1, left, top + i*gap + Inches(0.09),
                                     Pt(5), Pt(5))
        dot.fill.solid()
        dot.fill.fore_color.rgb = ACCENT
        dot.line.fill.background()
        txt(slide, item, left + Inches(0.18), top + i*gap,
            width - Inches(0.18), gap, size=size, color=color)


# ── Slide 1: Title ───────────────────────────────────────────────────────────
s = add_slide(); bg(s); accent_bar(s)
eyebrow(s, "RFP MC-2026-0417 · May 2026", top=Inches(1.4))
txt(s, "Modernizing the\nMeridian Dashboard", Inches(0.5), Inches(1.9),
    Inches(10), Inches(2.4), size=44, bold=True, color=WHITE)
txt(s, "Capabilities Presentation — Accenture",
    Inches(0.5), Inches(4.4), Inches(10), Inches(0.5), size=20, color=DIM)
divider(s, Inches(5.1))
txt(s, "Remediation  ·  New capability  ·  Test coverage  ·  Architecture handoff",
    Inches(0.5), Inches(5.3), Inches(12), Inches(0.5), size=13, color=DIM)

# ── Slide 2: Situation ───────────────────────────────────────────────────────
s = add_slide(); bg(s); accent_bar(s)
eyebrow(s, "The situation")
heading(s, "You have a working dashboard.\nIt's not fully working.", size=32)
body(s, "Three problems are blocking your operations team from getting full value.",
     Inches(0.5), Inches(2.6), Inches(12), Inches(0.4), size=16, color=MUTED)

cards = [
    ("REPORTS",    RGBColor(0x7f,0x1d,0x1d), RGBColor(0xfc,0xa5,0xa5),
     "Defects unresolved",
     "At least 8 known issues left at handoff. Filters, i18n gaps, inconsistent data patterns."),
    ("CAPABILITY", RGBColor(0x78,0x35,0x0f), AMBER,
     "No restocking view",
     "Your ops team can see stock levels but can't generate a purchase recommendation."),
    ("RISK",       RGBColor(0x7f,0x1d,0x1d), RGBColor(0xfc,0xa5,0xa5),
     "No test coverage",
     "IT has blocked changes. Without automated tests, every change is a manual risk assessment."),
]
for i, (tag, tf, tc, title, desc) in enumerate(cards):
    cx = Inches(0.5 + i * 4.28)
    box(s, cx, Inches(3.1), Inches(4.0), Inches(3.5), fill_color=CARD,
        border_color=RGBColor(0x33, 0x41, 0x55))
    tag_pill(s, tag, cx + Inches(0.2), Inches(3.3), tf, tc)
    txt(s, title, cx + Inches(0.2), Inches(3.72), Inches(3.6), Inches(0.4),
        size=15, bold=True, color=WHITE)
    txt(s, desc,  cx + Inches(0.2), Inches(4.2),  Inches(3.6), Inches(1.2),
        size=13, color=MUTED)

# ── Slide 3: Our read ────────────────────────────────────────────────────────
s = add_slide(); bg(s); accent_bar(s)
eyebrow(s, "Our read")
heading(s, "IT isn't the obstacle.\nIT is the unlock.", size=34)
body(s, "Most proposals treat testing as the last deliverable. We treat it as the first.\nUntil IT can approve changes with confidence, nothing you deliver is maintainable.",
     Inches(0.5), Inches(2.8), Inches(7.5), Inches(1.0), size=15)
body(s, "We establish test coverage in week two — before any defect fixes land.\nThat's how you get IT on your side instead of managing around them.",
     Inches(0.5), Inches(3.9), Inches(7.5), Inches(1.0), size=15)
box(s, Inches(8.8), Inches(2.3), Inches(4.0), Inches(3.8), fill_color=CARD,
    border_color=RGBColor(0x33,0x41,0x55))
tag_pill(s, "Sequencing", Inches(9.0), Inches(2.55),
         RGBColor(0x1e,0x3a,0x5f), RGBColor(0x93,0xc5,0xfd))
txt(s, "Tests → Reports → Restocking", Inches(9.0), Inches(3.0),
    Inches(3.6), Inches(0.5), size=14, bold=True, color=WHITE)
txt(s, "Each phase lands with coverage already in place. IT sees approvable changes from week 3 onward, not week 8.",
    Inches(9.0), Inches(3.55), Inches(3.6), Inches(1.8), size=13, color=MUTED)

# ── Slide 4: R1 ─────────────────────────────────────────────────────────────
s = add_slide(); bg(s); accent_bar(s)
eyebrow(s, "Requirement R1")
heading(s, "Reports remediation", size=36)
body(s, "No formal defect list exists — we audit everything, document findings, get ops team sign-off, then resolve all issues.",
     Inches(0.5), Inches(2.35), Inches(12), Inches(0.5), size=15)

box(s, Inches(0.5), Inches(3.1), Inches(5.8), Inches(3.5), fill_color=CARD,
    border_color=RGBColor(0x33,0x41,0x55))
tag_pill(s, "Approach", Inches(0.7), Inches(3.3),
         RGBColor(0x1e,0x3a,0x5f), RGBColor(0x93,0xc5,0xfd))
bullet_list(s, [
    "Full filter audit — Time Period, Warehouse, Category, Status",
    "API pattern consistency review",
    "i18n gap extraction to localisation layer",
    "Console error resolution",
], Inches(0.75), Inches(3.78), Inches(5.3))

box(s, Inches(6.9), Inches(3.1), Inches(5.9), Inches(3.5), fill_color=CARD,
    border_color=RGBColor(0x33,0x41,0x55))
tag_pill(s, "Deliverable", Inches(7.1), Inches(3.3),
         RGBColor(0x14,0x53,0x2d), GREEN)
txt(s, "Defect inventory shared with R. Tanaka's team for sign-off before we start fixing. No surprises.",
    Inches(7.1), Inches(3.78), Inches(5.5), Inches(1.2), size=14, color=MUTED)
txt(s, "✓  Complete by end of week 4",
    Inches(7.1), Inches(5.2), Inches(5.5), Inches(0.4), size=14, bold=True, color=GREEN)

# ── Slide 5: R2 ─────────────────────────────────────────────────────────────
s = add_slide(); bg(s); accent_bar(s)
eyebrow(s, "Requirement R2")
heading(s, "Restocking recommendations", size=36)
body(s, "Full optimization view — enter a budget ceiling, select warehouse scope, receive ranked purchase order recommendations.",
     Inches(0.5), Inches(2.35), Inches(12), Inches(0.5), size=15)

cols = [
    ("INPUTS",    RGBColor(0x1e,0x3a,0x5f), RGBColor(0x93,0xc5,0xfd),
     ["Current stock levels", "Demand forecast data", "Operator budget ceiling"]),
    ("LOGIC",     RGBColor(0x3b,0x07,0x64), PURPLE,
     ["Constrained allocation", "Maximize units within budget", "Server-side via /api/restocking"]),
    ("OUTPUT",    RGBColor(0x14,0x53,0x2d), GREEN,
     ["Ranked recommendation table", "SKU · stock · demand · qty · cost", "Per-warehouse scope"]),
]
for i, (tag, tf, tc, items) in enumerate(cols):
    cx = Inches(0.5 + i * 4.28)
    box(s, cx, Inches(3.1), Inches(4.0), Inches(3.6), fill_color=CARD,
        border_color=RGBColor(0x33,0x41,0x55))
    tag_pill(s, tag, cx + Inches(0.2), Inches(3.3), tf, tc)
    bullet_list(s, items, cx + Inches(0.2), Inches(3.78), Inches(3.6))

# ── Slide 6: R3 ─────────────────────────────────────────────────────────────
s = add_slide(); bg(s); accent_bar(s)
eyebrow(s, "Requirement R3")
heading(s, "Automated browser testing", size=36)
body(s, "Playwright end-to-end tests covering the four flows IT needs to approve future changes.",
     Inches(0.5), Inches(2.35), Inches(12), Inches(0.5), size=15)

box(s, Inches(0.5), Inches(3.1), Inches(5.8), Inches(3.5), fill_color=CARD,
    border_color=RGBColor(0x33,0x41,0x55))
tag_pill(s, "Covered flows", Inches(0.7), Inches(3.3),
         RGBColor(0x1e,0x3a,0x5f), RGBColor(0x93,0xc5,0xfd))
bullet_list(s, [
    "Dashboard — KPI tiles load with data",
    "Inventory — filter by warehouse and category",
    "Reports — all filters apply, no console errors",
    "Restocking — budget input, recommendations returned",
], Inches(0.75), Inches(3.78), Inches(5.3))

box(s, Inches(6.9), Inches(3.1), Inches(5.9), Inches(3.5), fill_color=CARD,
    border_color=RGBColor(0x33,0x41,0x55))
tag_pill(s, "Delivered as", Inches(7.1), Inches(3.3),
         RGBColor(0x14,0x53,0x2d), GREEN)
txt(s, "Tests live in the repository alongside application code — CI-ready from day one, not a one-time artifact.",
    Inches(7.1), Inches(3.78), Inches(5.5), Inches(1.2), size=14, color=MUTED)
txt(s, "✓  IT unblocked for Reports in week 4\n✓  Full coverage complete week 7",
    Inches(7.1), Inches(5.1), Inches(5.5), Inches(0.7), size=13, bold=True, color=GREEN)

# ── Slide 7: R4 ─────────────────────────────────────────────────────────────
s = add_slide(); bg(s); accent_bar(s)
eyebrow(s, "Requirement R4")
heading(s, "Architecture documentation", size=36)
body(s, "Produced during onboarding — before any code changes — so it reflects what was actually delivered, not what we modify.",
     Inches(0.5), Inches(2.35), Inches(12), Inches(0.5), size=15)
bullet_list(s, [
    "Component diagram: frontend views, API layer, backend, data layer",
    "Primary filter/query data flow",
    "Full API endpoint inventory",
    "Technical debt summary — what we inherited from the previous vendor",
], Inches(0.5), Inches(3.2), Inches(8.0), gap=Inches(0.38))
box(s, Inches(9.2), Inches(3.0), Inches(3.6), Inches(2.0), fill_color=CARD,
    border_color=RGBColor(0x33,0x41,0x55))
tag_pill(s, "Format", Inches(9.4), Inches(3.2),
         RGBColor(0x1e,0x3a,0x5f), RGBColor(0x93,0xc5,0xfd))
txt(s, "Interactive HTML diagram. No special tools required for Meridian IT.",
    Inches(9.4), Inches(3.65), Inches(3.2), Inches(1.0), size=13, color=MUTED)
txt(s, "Updated at engagement close and handed off to IT at the final review session.",
    Inches(0.5), Inches(5.3), Inches(12), Inches(0.5), size=13, color=DIM)

# ── Slide 8: Stretch ────────────────────────────────────────────────────────
s = add_slide(); bg(s); accent_bar(s)
eyebrow(s, "Desired items")
heading(s, "Stretch objectives at no additional cost", size=32)
body(s, "D1–D3 are in scope if timeline permits. Prioritized with your team after R1–R4 are delivered.",
     Inches(0.5), Inches(2.45), Inches(12), Inches(0.4), size=15)

stretch = [
    ("D1", RGBColor(0x78,0x35,0x0f), AMBER,
     "UI modernization",
     "Visual refresh to current standards — or your brand guide if you have one."),
    ("D2", RGBColor(0x78,0x35,0x0f), AMBER,
     "Full i18n",
     "Japanese-language support for Tokyo warehouse. Extends the i18n layer already partially in place."),
    ("D3", RGBColor(0x78,0x35,0x0f), AMBER,
     "Dark mode",
     "Operator-selectable theme for low-light floor stations. Prototyped on a branch — no risk to main."),
]
for i, (tag, tf, tc, title, desc) in enumerate(stretch):
    cx = Inches(0.5 + i * 4.28)
    box(s, cx, Inches(3.1), Inches(4.0), Inches(3.5), fill_color=CARD,
        border_color=RGBColor(0x33,0x41,0x55))
    tag_pill(s, tag, cx + Inches(0.2), Inches(3.3), tf, tc)
    txt(s, title, cx + Inches(0.2), Inches(3.72), Inches(3.6), Inches(0.4),
        size=15, bold=True, color=WHITE)
    txt(s, desc,  cx + Inches(0.2), Inches(4.2),  Inches(3.6), Inches(1.2),
        size=13, color=MUTED)
txt(s, "We don't commit to all three — but we commit to none of them blocking R1–R4.",
    Inches(0.5), Inches(6.8), Inches(12), Inches(0.4), size=13, color=DIM)

# ── Slide 9: Timeline ───────────────────────────────────────────────────────
s = add_slide(); bg(s); accent_bar(s)
eyebrow(s, "8-week engagement")
heading(s, "Timeline", size=36)

phases = [
    ("Weeks 1–2", "Phase 1: Foundation",
     "Architecture review\nTest infrastructure\nDefect audit",
     "R4 draft delivered"),
    ("Weeks 3–4", "Phase 2: Reports",
     "All defects resolved\nReports + Dashboard\ntest coverage",
     "R1 + partial R3"),
    ("Weeks 5–7", "Phase 3: Restocking",
     "API endpoint\nRestocking view\nInventory + Restocking tests",
     "R2 + R3 complete"),
    ("Week 8",    "Phase 4: Close",
     "R4 final\nStretch items\nIT handoff session",
     "Full acceptance"),
]
for i, (weeks, title, detail, milestone) in enumerate(phases):
    cx = Inches(0.5 + i * 3.2)
    box(s, cx, Inches(2.5), Inches(3.0), Inches(4.2), fill_color=CARD,
        border_color=RGBColor(0x33,0x41,0x55))
    txt(s, weeks, cx + Inches(0.2), Inches(2.68), Inches(2.7), Inches(0.3),
        size=10, bold=True, color=ACCENT)
    txt(s, title, cx + Inches(0.2), Inches(3.0), Inches(2.7), Inches(0.45),
        size=14, bold=True, color=WHITE)
    txt(s, detail, cx + Inches(0.2), Inches(3.5), Inches(2.7), Inches(1.2),
        size=12, color=MUTED)
    txt(s, f"→ {milestone}", cx + Inches(0.2), Inches(5.8), Inches(2.7), Inches(0.35),
        size=12, bold=True, color=GREEN)
txt(s, "Proposed start: May 18, 2026  ·  Completion: July 11, 2026",
    Inches(0.5), Inches(7.0), Inches(12), Inches(0.35), size=12, color=DIM)

# ── Slide 10: Pricing ───────────────────────────────────────────────────────
s = add_slide(); bg(s); accent_bar(s)
eyebrow(s, "Pricing")
heading(s, "Fixed fee. No surprises.", size=34)

# Left: price box
box(s, Inches(0.5), Inches(2.5), Inches(5.5), Inches(4.2),
    fill_color=RGBColor(0x1e,0x3a,0x5f),
    border_color=ACCENT)
txt(s, "Total fixed fee", Inches(0.75), Inches(2.75), Inches(5.0), Inches(0.35),
    size=12, color=MUTED)
txt(s, "$28,500", Inches(0.75), Inches(3.1), Inches(5.0), Inches(0.9),
    size=48, bold=True, color=WHITE)
txt(s, "Not-to-exceed · All R1–R4 deliverables included",
    Inches(0.75), Inches(4.05), Inches(5.0), Inches(0.4), size=13, color=MUTED)
divider(s, Inches(4.55))
txt(s, "30% on start  ·  40% at Phase 2  ·  30% at acceptance",
    Inches(0.75), Inches(4.65), Inches(5.0), Inches(0.4), size=12, color=DIM)

# Right: breakdown
rows = [
    ("R4 — Architecture documentation", "$3,000"),
    ("R1 — Reports remediation",         "$7,500"),
    ("R3 — Automated testing",           "$5,500"),
    ("R2 — Restocking recommendations",  "$10,500"),
    ("PM, handoff & buffer",             "$2,000"),
    ("Total",                            "$28,500"),
]
for i, (label, amount) in enumerate(rows):
    ry = Inches(2.5 + i * 0.62)
    is_total = label == "Total"
    border = ACCENT if is_total else RGBColor(0x33,0x41,0x55)
    box(s, Inches(6.5), ry, Inches(6.3), Inches(0.55),
        fill_color=CARD, border_color=border)
    txt(s, label, Inches(6.7), ry + Inches(0.12), Inches(4.5), Inches(0.35),
        size=13, bold=is_total, color=WHITE if is_total else MUTED)
    txt(s, amount, Inches(6.7), ry + Inches(0.12), Inches(5.9), Inches(0.35),
        size=13, bold=True, color=ACCENT, align=PP_ALIGN.RIGHT)

# ── Slide 11: Experience ─────────────────────────────────────────────────────
s = add_slide(); bg(s); accent_bar(s)
eyebrow(s, "Relevant experience")
heading(s, "We've done this before.", size=36)

refs = [
    ("2025 · VANTAGE LOGISTICS",  RGBColor(0x1e,0x3a,0x5f), RGBColor(0x93,0xc5,0xfd),
     "Warehouse Ops Dashboard",
     "Defect remediation + demand-driven reorder + Playwright coverage. Fixed-fee, on schedule. 4,000+ recommendations since go-live.",
     "→ Comparable to R1, R2, R3"),
    ("2024 · CORSTONE MFG",       RGBColor(0x3b,0x07,0x64), PURPLE,
     "Inventory Control System",
     "Inherited partial build. Architecture review, i18n for Japanese facility, IT handoff documentation at close.",
     "→ Comparable to R4, D2"),
    ("2024 · REDWOOD PARTS",      RGBColor(0x14,0x53,0x2d), GREEN,
     "Analytics Platform Testing",
     "Built Playwright coverage from zero. Eight flows, CI-gated. IT sign-off within two weeks of delivery.",
     "→ Comparable to R3"),
]
for i, (tag, tf, tc, title, desc, comp) in enumerate(refs):
    cx = Inches(0.5 + i * 4.28)
    box(s, cx, Inches(2.5), Inches(4.0), Inches(4.3), fill_color=CARD,
        border_color=RGBColor(0x33,0x41,0x55))
    tag_pill(s, tag, cx + Inches(0.2), Inches(2.72), tf, tc)
    txt(s, title, cx + Inches(0.2), Inches(3.12), Inches(3.6), Inches(0.4),
        size=14, bold=True, color=WHITE)
    txt(s, desc,  cx + Inches(0.2), Inches(3.6),  Inches(3.6), Inches(1.5),
        size=12, color=MUTED)
    txt(s, comp,  cx + Inches(0.2), Inches(5.8),  Inches(3.6), Inches(0.35),
        size=12, bold=True, color=ACCENT2)

# ── Slide 12: Close ──────────────────────────────────────────────────────────
s = add_slide(); bg(s); accent_bar(s)
eyebrow(s, "Next steps")
heading(s, "Ready to start\nweek of May 18.", size=38)
body(s, "Your ops team has been working around a broken Reports module and a missing\nRestocking capability for long enough. We can close both gaps — and leave you\nwith the test infrastructure to keep them closed.",
     Inches(0.5), Inches(3.2), Inches(8.0), Inches(1.4), size=16)

box(s, Inches(0.5), Inches(4.9), Inches(3.8), Inches(1.4), fill_color=CARD,
    border_color=RGBColor(0x33,0x41,0x55))
tag_pill(s, "Fixed fee", Inches(0.7), Inches(5.1),
         RGBColor(0x1e,0x3a,0x5f), RGBColor(0x93,0xc5,0xfd))
txt(s, "$28,500", Inches(0.7), Inches(5.48), Inches(3.4), Inches(0.45),
    size=18, bold=True, color=WHITE)
txt(s, "All four required deliverables.", Inches(0.7), Inches(5.9),
    Inches(3.4), Inches(0.3), size=12, color=MUTED)

box(s, Inches(4.6), Inches(4.9), Inches(3.8), Inches(1.4), fill_color=CARD,
    border_color=RGBColor(0x33,0x41,0x55))
tag_pill(s, "Timeline", Inches(4.8), Inches(5.1),
         RGBColor(0x14,0x53,0x2d), GREEN)
txt(s, "8 weeks", Inches(4.8), Inches(5.48), Inches(3.4), Inches(0.45),
    size=18, bold=True, color=WHITE)
txt(s, "IT unblocked week 4. Delivery July 11.", Inches(4.8), Inches(5.9),
    Inches(3.4), Inches(0.3), size=12, color=MUTED)

txt(s, "Jakub Raczek · Accenture  ·  procurement@meridiancomponents.example",
    Inches(0.5), Inches(7.0), Inches(12), Inches(0.35), size=12, color=DIM)


out = r"c:\CProjektyGIT\Claude\meridian-workshop\proposal\capabilities-deck.pptx"
prs.save(out)
print(f"Saved: {out}")
