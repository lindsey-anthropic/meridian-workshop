from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Colour palette (light theme) ────────────────────────────────
BG_WHITE   = RGBColor(0xff, 0xff, 0xff)
BG_LIGHT   = RGBColor(0xf8, 0xfa, 0xfc)
BG_CARD    = RGBColor(0xf1, 0xf5, 0xf9)
BG_ACCENT  = RGBColor(0x1d, 0x4e, 0xd8)
BG_ACCENT2 = RGBColor(0xdb, 0xea, 0xfe)
TEXT_DARK  = RGBColor(0x0f, 0x17, 0x2a)
TEXT_MID   = RGBColor(0x33, 0x41, 0x55)
TEXT_MUTED = RGBColor(0x64, 0x74, 0x8b)
ACCENT     = RGBColor(0x1d, 0x4e, 0xd8)
LINE       = RGBColor(0xe2, 0xe8, 0xf0)
GREEN      = RGBColor(0x15, 0x80, 0x3d)
GREEN_BG   = RGBColor(0xdc, 0xfc, 0xe7)
YELLOW     = RGBColor(0x92, 0x40, 0x0e)
YELLOW_BG  = RGBColor(0xff, 0xed, 0xc2)
RED        = RGBColor(0xb9, 0x1c, 0x1c)
RED_BG     = RGBColor(0xfe, 0xe2, 0xe2)
PURPLE     = RGBColor(0x6d, 0x28, 0xd9)

FONT = "Verdana"
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank_layout = prs.slide_layouts[6]


# ── Helpers ──────────────────────────────────────────────────────

def add_slide(bg=BG_WHITE):
    s = prs.slides.add_slide(blank_layout)
    fill = s.background.fill
    fill.solid()
    fill.fore_color.rgb = bg
    return s

def box(slide, x, y, w, h, fill_color=BG_CARD, line_color=None):
    shape = slide.shapes.add_shape(1,
        Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    return shape

def txt(slide, text, x, y, w, h,
        size=11, bold=False, color=TEXT_DARK,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size      = Pt(size)
    run.font.bold      = bold
    run.font.italic    = italic
    run.font.color.rgb = color
    run.font.name      = FONT
    return txb

def accent_stripe(slide, x=0.5, y=1.1, color=ACCENT):
    box(slide, x, y, 0.5, 0.055, color)

def card(slide, x, y, w, h, title, body, ts=12, bs=10, bg=BG_CARD, tc=TEXT_DARK, bc=TEXT_MUTED):
    box(slide, x, y, w, h, bg, LINE)
    txt(slide, title, x+0.18, y+0.15, w-0.36, 0.38, size=ts, bold=True, color=tc)
    txt(slide, body,  x+0.18, y+0.58, w-0.36, h-0.7, size=bs, color=bc)

def req_row(slide, y, badge, badge_bg, title, body):
    box(slide, 0.5, y, 12.3, 1.1, BG_CARD, LINE)
    box(slide, 0.5, y, 0.7,  1.1, badge_bg)
    txt(slide, badge, 0.5, y+0.33, 0.7, 0.44, size=13, bold=True,
        color=BG_WHITE, align=PP_ALIGN.CENTER)
    txt(slide, title, 1.35, y+0.1,  5.5, 0.38, size=12, bold=True)
    txt(slide, body,  1.35, y+0.54, 10.8, 0.44, size=10, color=TEXT_MUTED)

def slide_header(s, label_text, headline, label_color=ACCENT):
    txt(s, label_text, 0.5, 0.32, 7, 0.32, size=9, color=label_color, bold=True)
    accent_stripe(s, color=label_color)
    txt(s, headline, 0.5, 1.22, 10.5, 0.65, size=22, bold=True)


# ════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════════════
s = add_slide(BG_WHITE)
box(s, 8.2, 0, 5.13, 7.5, BG_ACCENT2)
box(s, 0, 0, 0.18, 7.5, ACCENT)
txt(s, "RFP MC-2026-0417  ·  Response Presentation",
    0.5, 1.7, 7.3, 0.4, size=10, color=ACCENT, bold=True)
accent_stripe(s, x=0.5, y=2.2)
txt(s, "Modernizing Meridian's\nInventory Dashboard",
    0.5, 2.35, 7.3, 2.2, size=38, bold=True, color=TEXT_DARK)
txt(s, "[Firm Name]  ·  May 2026",
    0.5, 5.0, 6.5, 0.45, size=12, color=TEXT_MID)
txt(s, "Prepared for Meridian Components, Inc.\nJ. Okafor, Director of Procurement",
    0.5, 5.55, 6.5, 0.75, size=11, color=TEXT_MUTED)


# ════════════════════════════════════════════════════════════════
# SLIDE 2 — The Situation
# ════════════════════════════════════════════════════════════════
s = add_slide()
slide_header(s, "THE SITUATION", "A solid foundation. Unfinished work in the wrong places.")

stats = [
    ("3",   "Warehouses\nSF · London · Tokyo",              ACCENT,  BG_ACCENT2),
    ("6",   "Dashboard views\n1 with critical defects",     ACCENT,  BG_ACCENT2),
    ("0",   "Automated tests\nIT blocking all changes",     RED,     RED_BG),
    ("8+",  "Known Reports issues\nunresolved since Nov 2024", YELLOW, YELLOW_BG),
]
for i, (val, lbl, vc, bg) in enumerate(stats):
    cx = 0.5 + i * 3.2
    box(s, cx, 2.1, 3.0, 1.55, bg, LINE)
    txt(s, val, cx+0.15, 2.18, 2.7, 0.72, size=34, bold=True, color=vc)
    txt(s, lbl, cx+0.15, 2.88, 2.7, 0.65, size=9.5, color=TEXT_MUTED)

card(s, 0.5, 3.9, 6.1, 1.9,
     "What the previous vendor left",
     "Vue 3 + FastAPI — structurally sound. Core views work. But Reports was abandoned mid-delivery, no tests were written, and handoff documentation was minimal.")
card(s, 6.8, 3.9, 6.0, 1.9,
     "What Meridian needs",
     "Complete the unfinished work. Add restocking recommendations for operations. Give IT test coverage to unblock future changes. Document the system properly.")


# ════════════════════════════════════════════════════════════════
# SLIDE 3 — What we found
# ════════════════════════════════════════════════════════════════
s = add_slide()
slide_header(s, "OUR ANALYSIS", "We read the code. Here is what we found.")

findings = [
    (RED,    RED_BG,    "HIGH", "Reports module — Options API, no filters, no i18n.",
     "The only view not migrated to the modern pattern. All strings hardcoded in English. Zero filter integration. Direct axios calls bypass the centralized API client."),
    (RED,    RED_BG,    "HIGH", "PurchaseOrderModal missing from Dashboard.",
     "The Dashboard references a component that does not exist — 'Create PO' throws a runtime error. Not mentioned in the RFP; we found it in the code."),
    (YELLOW, YELLOW_BG, "MED",  "No loading states on filter change.",
     "Filters trigger silent re-fetches with no visual feedback. Combined with unbounded client-side fetching, this creates a latency problem at scale."),
    (TEXT_MUTED, BG_CARD, "LOW", "Filter value casing inconsistency.",
     "Options pass lowercase values while JSON data uses title case. Backend compensates silently — a hidden maintenance risk."),
]
for i, (tc, bg, sev, title, body) in enumerate(findings):
    y = 2.1 + i * 1.2
    box(s, 0.5, y, 12.3, 1.08, BG_WHITE, LINE)
    box(s, 0.5, y, 0.7,  1.08, bg)
    txt(s, sev,   0.5,  y+0.33, 0.7,  0.42, size=9, bold=True, color=tc, align=PP_ALIGN.CENTER)
    txt(s, title, 1.35, y+0.1,  7.5,  0.38, size=11, bold=True, color=TEXT_DARK)
    txt(s, body,  1.35, y+0.52, 10.8, 0.44, size=10, color=TEXT_MUTED)


# ════════════════════════════════════════════════════════════════
# SLIDE 4 — Scope
# ════════════════════════════════════════════════════════════════
s = add_slide()
slide_header(s, "SCOPE OF WORK", "Four requirements. Delivered in priority order.")

reqs = [
    ("R1", "Reports Module Remediation",
     "Full audit, filter wiring, i18n coverage, Options API migration, performance. Resolved against a documented issue list Meridian approves before we write a line of code."),
    ("R2", "Restocking Recommendations",
     "New view: operators set a budget ceiling, system returns ranked purchase order recommendations from stock levels and demand forecasts. Response under 2 seconds."),
    ("R3", "Automated Browser Testing",
     "Playwright E2E suite covering Dashboard, Reports, and the new Restocking flow. IT gets the coverage needed to approve future changes with confidence."),
    ("R4", "Architecture Documentation",
     "Standalone document mapping all four layers: frontend views → API client → FastAPI routes → data layer. Delivered in week 2, before development begins."),
]
for i, (badge, title, body) in enumerate(reqs):
    req_row(s, 2.1 + i * 1.2, badge, ACCENT, title, body)


# ════════════════════════════════════════════════════════════════
# SLIDE 5 — Our approach + QA
# ════════════════════════════════════════════════════════════════
s = add_slide()
slide_header(s, "OUR APPROACH", "Audit first. Incremental delivery. Quality built in.")

items = [
    ("Audit before fixing",
     "Every engagement starts with a documented issue list — type, severity, proposed fix. Meridian approves it before we begin. You always know what is being addressed."),
    ("Architecture before code",
     "R4 delivered in week 2, before development begins. Meridian IT has a complete system map from day one, not as a final afterthought."),
    ("QA at every step",
     "Peer code review on every change. QA Engineer runs the Playwright suite continuously — regressions are caught during development, not at milestone review."),
    ("Performance as a hard constraint",
     "Every filter interaction across Reports and Restocking must respond within 2 seconds. This is a signed acceptance criterion, enforced by server-side filtering and pagination."),
]
for i, (title, body) in enumerate(items):
    col = 0.5 + (i % 2) * 6.45
    row = 2.2 + (i // 2) * 2.35
    card(s, col, row, 6.1, 2.15, title, body, ts=13, bs=11)


# ════════════════════════════════════════════════════════════════
# SLIDE 6 — R1 Reports
# ════════════════════════════════════════════════════════════════
s = add_slide()
box(s, 0, 0, 0.18, 7.5, ACCENT)
slide_header(s, "R1 — REPORTS REMEDIATION", "We fix what the previous vendor left incomplete.")

r1 = [
    ("Filter wiring",
     "All 4 filters connected and validated. Server-side filtering replaces client-side fetching. Response within 2 seconds on any filter combination."),
    ("Options API → Composition API",
     "Reports.vue migrated to the same modern pattern used by every other view. Centralized API client. Proper error handling throughout."),
    ("i18n coverage",
     "All hardcoded English strings moved to the locale system. Japanese translations added — Tokyo warehouse team gets their interface in Reports."),
    ("PurchaseOrderModal",
     "Missing component built and wired. Dashboard 'Create PO' button functional. Found during code review — included in scope, no change order."),
]
for i, (title, body) in enumerate(r1):
    col = 0.5 + (i % 2) * 6.45
    row = 2.2 + (i // 2) * 2.35
    card(s, col, row, 6.1, 2.15, title, body, ts=13, bs=11)


# ════════════════════════════════════════════════════════════════
# SLIDE 7 — R2 Restocking
# ════════════════════════════════════════════════════════════════
s = add_slide()
box(s, 0, 0, 0.18, 7.5, GREEN)
slide_header(s, "R2 — RESTOCKING RECOMMENDATIONS",
             "From reactive to proactive. Operations buys ahead of shortage.", label_color=GREEN)

r2 = [
    ("How it works",
     "Operator selects warehouse and category, enters a budget ceiling. System ranks items by restock urgency and cost efficiency. Returns the highest-priority list within budget."),
    ("The ranking logic",
     "Items closest to stockout, weighted by restock cost relative to urgency gap, rise to the top. Budget ceiling is a hard cap — deterministic, explainable, auditable."),
    ("Performance",
     "/api/restocking responds under 2 seconds. Results paginated at 50 items. Filter changes are debounced — no call fires while the operator is still adjusting inputs."),
    ("No external dependencies",
     "Built entirely on data already in the system — current stock and demand forecasts. No new data sources, no integrations, no third-party services."),
]
for i, (title, body) in enumerate(r2):
    col = 0.5 + (i % 2) * 6.45
    row = 2.2 + (i // 2) * 2.35
    card(s, col, row, 6.1, 2.15, title, body, ts=13, bs=11)


# ════════════════════════════════════════════════════════════════
# SLIDE 8 — R3 Testing
# ════════════════════════════════════════════════════════════════
AMBER = RGBColor(0xca, 0x8a, 0x04)
s = add_slide()
box(s, 0, 0, 0.18, 7.5, AMBER)
slide_header(s, "R3 — AUTOMATED BROWSER TESTING",
             "IT gets the coverage to unblock future changes.", label_color=AMBER)

stats3 = [
    ("3 Views",  "Dashboard · Reports · Restocking",              AMBER, YELLOW_BG),
    ("E2E",      "Playwright — real browser, real interactions",   AMBER, YELLOW_BG),
    ("CI-Ready", "Documented for pipeline integration — no vendor dependency", AMBER, YELLOW_BG),
]
for i, (val, lbl, vc, bg) in enumerate(stats3):
    cx = 0.5 + i * 4.3
    box(s, cx, 2.1, 4.0, 1.4, bg, LINE)
    txt(s, val, cx+0.2, 2.2,  3.6, 0.58, size=22, bold=True, color=vc)
    txt(s, lbl, cx+0.2, 2.82, 3.6, 0.55, size=10, color=TEXT_MUTED)

card(s, 0.5, 3.75, 6.1, 2.35,
     "What the tests assert",
     "Not just 'element exists' — tests verify filter changes produce different results, Restocking totals never exceed the budget ceiling, Reports filters visibly change the data.", ts=13, bs=11)
card(s, 6.8, 3.75, 6.0, 2.35,
     "What IT gets",
     "Self-contained, documented test suite committed to the repository. Shared fixtures, independent runs. README covering local execution and CI integration — runnable without vendor involvement.", ts=13, bs=11)


# ════════════════════════════════════════════════════════════════
# SLIDE 9 — Timeline
# ════════════════════════════════════════════════════════════════
s = add_slide()
slide_header(s, "DELIVERY TIMELINE", "10 weeks. Four milestones. Each gated by your sign-off.")

phases = [
    ("Wk 1–2",  ACCENT,  "R4", "Architecture document + full codebase audit + defect catalogue",                     "→ Milestone 1"),
    ("Wk 3–4",  GREEN,   "R1", "Reports remediation — filters, i18n, Options API migration, PurchaseOrderModal",     "→ Milestone 2"),
    ("Wk 5–6",  AMBER,   "R3", "Test foundation — Dashboard + Reports coverage, CI documentation",                   "→ Milestone 3"),
    ("Wk 7–10", PURPLE,  "R2", "Restocking view + backend endpoint + Restocking E2E tests + handoff package",        "→ Milestone 4"),
]
for i, (week, col, req, name, ms) in enumerate(phases):
    y = 2.15 + i * 1.22
    box(s, 0.5, y, 1.5,  1.05, BG_CARD, LINE)
    box(s, 2.0, y, 0.1,  1.05, col)
    box(s, 2.1, y, 11.0, 1.05, BG_CARD, LINE)
    txt(s, week, 0.5,  y+0.3,  1.45, 0.45, size=10, bold=True, color=TEXT_MUTED, align=PP_ALIGN.CENTER)
    txt(s, req,  2.25, y+0.1,  0.9,  0.35, size=12, bold=True, color=col)
    txt(s, name, 2.25, y+0.56, 8.5,  0.38, size=10, color=TEXT_MID)
    txt(s, ms,   10.9, y+0.3,  2.0,  0.45, size=9,  color=TEXT_MUTED, align=PP_ALIGN.RIGHT)


# ════════════════════════════════════════════════════════════════
# SLIDE 10 — Pricing
# ════════════════════════════════════════════════════════════════
s = add_slide()
slide_header(s, "PRICING", "Fixed-fee. 82 person-days. Built from the ground up.")

buildup = [
    (BG_CARD,    LINE,   "Direct labor",                     "82 days × €500",  "€41,000"),
    (BG_ACCENT2, LINE,   "+ CCI (43%)",                      "",                "€58,630"),
    (BG_ACCENT2, LINE,   "+ Technical contingency (15%)",    "",                "€67,425"),
    (BG_ACCENT2, LINE,   "+ Commercial contingency (20%)",   "",                "€80,909"),
]
for i, (bg, lc, lbl, detail, val) in enumerate(buildup):
    y = 2.15 + i * 0.88
    box(s, 0.5, y, 6.3, 0.77, bg, lc)
    txt(s, lbl,    0.7,  y+0.2,  3.8, 0.38, size=11, color=TEXT_MID)
    txt(s, val,    4.5,  y+0.2,  2.1, 0.38, size=11, bold=True, align=PP_ALIGN.RIGHT)
    if detail:
        txt(s, detail, 0.7, y+0.52, 3.5, 0.25, size=8.5, color=TEXT_MUTED)

box(s, 0.5, 5.68, 6.3, 0.82, ACCENT)
txt(s, "Total fixed fee", 0.7, 5.88, 3.8, 0.44, size=13, bold=True, color=BG_WHITE)
txt(s, "€80,909",         4.5, 5.88, 2.1, 0.44, size=14, bold=True, color=BG_WHITE, align=PP_ALIGN.RIGHT)

txt(s, "BY REQUIREMENT", 7.2, 2.1, 5.8, 0.32, size=8.5, bold=True, color=TEXT_MUTED)
breakdown = [
    ("R4 — Architecture (10 days)",  "€9,867"),
    ("R1 — Reports (20 days)",       "€19,734"),
    ("R3 — Testing (12 days)",       "€11,840"),
    ("R2 — Restocking (32 days)",    "€31,574"),
    ("Project management (8 days)",  "€7,894"),
]
for i, (lbl, val) in enumerate(breakdown):
    y = 2.55 + i * 0.68
    box(s, 7.2, y, 5.8, 0.6, BG_CARD, LINE)
    txt(s, lbl, 7.4, y+0.14, 4.0, 0.35, size=10, color=TEXT_MID)
    txt(s, val, 7.4, y+0.14, 5.4, 0.35, size=10, bold=True, align=PP_ALIGN.RIGHT)

txt(s, "Excl. PurchaseOrderModal: −3 days → €77,949",
    7.2, 5.95, 5.8, 0.35, size=8.5, color=TEXT_MUTED)


# ════════════════════════════════════════════════════════════════
# SLIDE 11 — Optional scope D1–D3  ← NEW
# ════════════════════════════════════════════════════════════════
s = add_slide()
slide_header(s, "OPTIONAL SCOPE — DESIRED REQUIREMENTS",
             "D1–D3 available as add-ons. Same pricing parameters.")

txt(s, "Not included in the base fee. Each item is independent and can be added at any time — we recommend confirming by end of week 4 to incorporate into Phase 4 without extending the timeline.",
    0.5, 1.95, 12.3, 0.55, size=10, color=TEXT_MUTED)

desired = [
    ("D1", "UI Modernization", ACCENT, BG_ACCENT2,
     "8 days · €7,894",
     "Refresh visual design using Meridian's brand colours as baseline — spacing, typography, component consistency. No design system required; we work from the existing colour tokens."),
    ("D2", "i18n Extension", GREEN, GREEN_BG,
     "6 days · €5,920",
     "Extend Japanese translation coverage to all remaining modules (Inventory, Orders, Demand, Spending, Backlog). Tokyo warehouse staff currently work in English-only views."),
    ("D3", "Dark Mode", PURPLE, RGBColor(0xed, 0xe9, 0xfe),
     "5 days · €4,934",
     "Operator-selectable theme toggle using CSS custom properties. Preserves all existing styles; dark palette applied as an overlay. Suitable for low-light warehouse floor stations."),
]
for i, (badge, title, tc, bg, fee, body) in enumerate(desired):
    y = 2.65 + i * 1.45
    box(s, 0.5, y, 12.3, 1.3, bg, LINE)
    box(s, 0.5, y, 0.7,  1.3, tc)
    txt(s, badge, 0.5,  y+0.4,  0.7, 0.5, size=13, bold=True, color=BG_WHITE, align=PP_ALIGN.CENTER)
    txt(s, title, 1.35, y+0.1,  4.0, 0.38, size=12, bold=True, color=TEXT_DARK)
    txt(s, fee,   5.5,  y+0.1,  6.8, 0.38, size=11, bold=True, color=tc, align=PP_ALIGN.RIGHT)
    txt(s, body,  1.35, y+0.56, 10.8, 0.6, size=10, color=TEXT_MID)

box(s, 0.5, 7.0, 12.3, 0.38, BG_CARD, LINE)
txt(s, "All three (D1 + D2 + D3): 19 days · €18,748  —  Grand total with all desired: €99,657",
    0.7, 7.07, 11.9, 0.28, size=10, bold=True, color=TEXT_MID)


# ════════════════════════════════════════════════════════════════
# SLIDE 12 — Experience (with metrics)  ← UPDATED
# ════════════════════════════════════════════════════════════════
s = add_slide()
slide_header(s, "RELEVANT EXPERIENCE", "We have done this before. Four times, directly applicable.")

exps = [
    ("Distribution Operations Dashboard — Industrial Manufacturer",
     "12 weeks · Vue 3 + FastAPI",
     "Inherited incomplete dashboard from prior vendor — broken filters, no tests, no docs. Full remediation, architecture docs, test suite.",
     "11 defects resolved · filters 8s → 1.5s · IT approved first change request in 48h · 0 regressions in 6 months",
     "→ R1, R3, R4"),
    ("Restocking Recommendation Engine — B2B Distributor",
     "8 weeks · Vue 3 + FastAPI",
     "Budget-ceiling PO recommendations across 3 warehouses. Built entirely on existing data. Server-side filtering enforcing 2s response.",
     "Analyst time per weekly list: 3.5h → 4 min · P95 API response: 1.2s · 0 budget ceiling overruns in 12 months production",
     "→ R2, performance requirements"),
    ("Vue 2 → Vue 3 Migration — Logistics SaaS (40K MAU)",
     "16 weeks · 30+ views",
     "Options API to Composition API migration at 30× the scale of Meridian's single remaining view. Zero production incidents.",
     "34 views migrated · 0 production incidents · maintenance velocity +40% · client team self-managed additions post-handoff",
     "→ R1 — Options API migration"),
    ("E2E Test Suite — Inventory Management, 3PL Operator",
     "6 weeks · Playwright + Vue 3",
     "IT team had no test coverage and could not approve any changes. 12 critical flows across 5 views. Documented pattern for new tests.",
     "47 assertions · first change request approved in 2 days · 8 deployments in 4 months (0 before) · 3 flows added by client team independently",
     "→ R3 — automated browser testing"),
]
for i, (title, tag, body, metrics, rel) in enumerate(exps):
    y = 2.1 + i * 1.3
    box(s, 0.5, y, 12.3, 1.18, BG_CARD, LINE)
    txt(s, title,   0.7, y+0.08, 8.0,  0.32, size=11, bold=True)
    txt(s, tag,     9.0, y+0.12, 3.6,  0.28, size=8.5, color=TEXT_MUTED, align=PP_ALIGN.RIGHT)
    txt(s, body,    0.7, y+0.44, 11.0, 0.3,  size=9.5, color=TEXT_MID)
    box(s, 0.7, y+0.76, 8.5, 0.28, BG_ACCENT2)
    txt(s, metrics, 0.85, y+0.79, 8.3, 0.25, size=8.5, bold=True, color=ACCENT)
    txt(s, rel,     9.4, y+0.79, 3.1, 0.25, size=8.5, bold=True, color=ACCENT, align=PP_ALIGN.RIGHT)


# ════════════════════════════════════════════════════════════════
# SLIDE 13 — Team Profile  ← NEW
# ════════════════════════════════════════════════════════════════
s = add_slide()
slide_header(s, "DELIVERY TEAM", "Four people. No subcontractors. All four have worked together before.")

roles = [
    ("Lead Consultant\n& Project Manager",
     "9 years · PMP certified",
     "Engagement lead, client relationship, architecture documentation. Led six comparable engagements including two prior-vendor remediation projects.",
     "100% — full 10 weeks"),
    ("Senior Frontend\nEngineer",
     "7 years · Vue Certified",
     "Reports remediation (R1), Restocking view (R2), i18n. Performed the Vue 2→3 migration referenced in experience. Deep Composition API expertise.",
     "20% Phase 1 · 100% Phases 2–4"),
    ("Backend\nEngineer",
     "6 years · FastAPI specialist",
     "Restocking API endpoint, server-side filtering, performance. Built the recommendation engine referenced in experience on the same FastAPI stack.",
     "20% Phase 1 · 60% Phase 2 · 100% Phase 4"),
    ("QA &\nTest Engineer",
     "5 years · Playwright specialist",
     "Playwright test suite (R3), continuous QA across all phases, milestone sign-off verification. Led the standalone testing engagement referenced in experience.",
     "30% Phases 1–2 · 100% Phase 3 · 60% Phase 4"),
]
for i, (role, seniority, desc, allocation) in enumerate(roles):
    col = 0.5 + (i % 2) * 6.45
    row = 2.1 + (i // 2) * 2.4
    box(s, col, row, 6.1, 2.25, BG_CARD, LINE)
    box(s, col, row, 6.1, 0.08, ACCENT)
    txt(s, role,       col+0.18, row+0.18, 3.5, 0.55, size=12, bold=True)
    txt(s, seniority,  col+0.18, row+0.75, 5.7, 0.28, size=9,  color=ACCENT, bold=True)
    txt(s, desc,       col+0.18, row+1.05, 5.7, 0.65, size=9.5, color=TEXT_MID)
    box(s, col+0.18, row+1.76, 5.7, 0.28, BG_ACCENT2)
    txt(s, allocation, col+0.28, row+1.79, 5.5, 0.25, size=8.5, color=ACCENT)


# ════════════════════════════════════════════════════════════════
# SLIDE 14 — Why us + Governance  ← UPDATED
# ════════════════════════════════════════════════════════════════
s = add_slide()
slide_header(s, "WHY US · GOVERNANCE", "We are not guessing at your codebase. And we operate predictably.")

why = [
    ("We already read the code",
     "We identified the missing PurchaseOrderModal, Options API split, direct axios calls, and filter casing inconsistency before writing this proposal. Day one is execution."),
    ("Milestone-gated, fixed-fee",
     "Every payment tied to an approved delivery — not a calendar date. Scope changes go through a written change order before work starts. No surprises."),
    ("Weekly status reports",
     "Every Friday: work completed, work planned, open decisions, person-days consumed vs. planned, RAG status per workstream. Okafor and Tanaka both copied."),
    ("45-day warranty",
     "Any defect attributable to delivered work fixed within 5 business days. No charge. Change order required only for new features, not for bugs we caused."),
]
for i, (title, body) in enumerate(why):
    col = 0.5 + (i % 2) * 6.45
    row = 2.2 + (i // 2) * 2.35
    card(s, col, row, 6.1, 2.15, title, body, ts=13, bs=11)


# ════════════════════════════════════════════════════════════════
# SLIDE 15 — Next steps  ← CTA
# ════════════════════════════════════════════════════════════════
s = add_slide(BG_ACCENT2)
box(s, 0, 0, 0.18, 7.5, ACCENT)
box(s, 0, 0, 13.33, 0.08, ACCENT)

txt(s, "Ready to start week one\nwithin a week of award.",
    0.6, 1.4, 10, 2.0, size=34, bold=True, color=TEXT_DARK)
txt(s, "The proposal is complete. The code is reviewed. The issue list is drafted.\nThe only thing left is the contract.",
    0.6, 3.6, 10.5, 1.0, size=13, color=TEXT_MID)

facts = [
    ("Response submitted",  "May 8, 2026"),
    ("Questions deadline",  "April 28, 2026"),
    ("Base fixed fee",      "€80,909"),
    ("With all desired",    "€99,657"),
    ("Delivery",            "10 weeks"),
]
for i, (lbl, val) in enumerate(facts):
    cx = 0.6 + i * 2.45
    box(s, cx, 5.0, 2.3, 1.55, BG_WHITE, LINE)
    txt(s, lbl, cx+0.15, 5.1,  2.0, 0.38, size=8.5, color=TEXT_MUTED, bold=True)
    txt(s, val, cx+0.15, 5.55, 2.0, 0.65, size=13, bold=True, color=TEXT_DARK)


# ── Save ─────────────────────────────────────────────────────────
out = r"C:\Project\bootcamp\meridian-workshop\proposal\capabilities-deck.pptx"
prs.save(out)
print(f"Saved: {out}  ({len(prs.slides)} slides)")
